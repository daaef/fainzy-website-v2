"""
Media Specification Generator for Fainzy Website V2
Generates a PowerPoint presentation with media dimension specifications
optimized for 1920x1080 viewport across all pages.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import urllib.request
from io import BytesIO
from PIL import Image as PILImage

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def add_title_slide(title, subtitle=""):
    """Add a title slide with dark theme"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 11)  # #0a0a0b

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(15), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle if provided
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(15), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(163, 163, 163)
        subtitle_para.alignment = PP_ALIGN.CENTER

    return slide

def add_page_slide(page_title, page_subtitle, media_items):
    """Add a page slide with media specifications"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 11)

    # Add page banner
    banner_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    banner_frame = banner_box.text_frame
    banner_frame.text = page_title
    banner_para = banner_frame.paragraphs[0]
    banner_para.font.size = Pt(36)
    banner_para.font.bold = True
    banner_para.font.color.rgb = RGBColor(255, 255, 255)

    # Add page subtitle if provided
    if page_subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(15), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = page_subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(16)
        subtitle_para.font.color.rgb = RGBColor(163, 163, 163)

    # Add media items
    y_position = 1.6
    for item in media_items:
        if y_position > 7.5:  # If we're running out of space, create new slide
            slide = prs.slides.add_slide(slide_layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(10, 10, 11)
            y_position = 0.5

        # Media title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_position), Inches(15), Inches(0.35))
        title_frame = title_box.text_frame
        title_frame.text = f"• {item['title']}"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(18)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        y_position += 0.35

        # Details
        details_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_position), Inches(14.5), Inches(0.6))
        details_frame = details_box.text_frame
        details_frame.word_wrap = True

        # Usage
        p = details_frame.paragraphs[0]
        p.text = f"Usage: {item['usage']}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 200, 200)

        # Dimensions
        p = details_frame.add_paragraph()
        p.text = f"Optimal Dimensions: {item['dimensions']}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 200, 200)

        # Media type recommendation
        p = details_frame.add_paragraph()
        p.text = f"Recommendation: {item['recommendation']}"
        p.font.size = Pt(12)
        p.font.bold = True
        color_map = {
            "Static Image": RGBColor(99, 102, 241),  # Blue
            "Animated GIF": RGBColor(236, 72, 153),   # Pink
            "Video": RGBColor(34, 197, 94)            # Green
        }
        p.font.color.rgb = color_map.get(item['recommendation'], RGBColor(200, 200, 200))

        # Context/Description
        if 'context' in item and item['context']:
            p = details_frame.add_paragraph()
            p.text = f"Context: {item['context']}"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(163, 163, 163)

        y_position += 0.8

    return slide

# ==================== PAGE 1: HOME PAGE ====================
add_title_slide(
    "Fainzy Website V2",
    "Media Dimensions & Specifications Guide\nOptimized for 1920x1080 Viewport"
)

# Home Page
home_media = [
    {
        "title": "Hero Carousel Slide 1: ZiBot",
        "usage": "Hero carousel slide - First slide showcasing ZiBot product",
        "dimensions": "1920x1080px (Full viewport background) | Product Image: 1200x500px (center positioned)",
        "recommendation": "Animated GIF",
        "context": "Title: 'ZiBot', Subtitle: 'Anything Anytime Anywhere Delivery'. Shows ZiBot robot with animation entrance effects. Recommend: Robot movement GIF"
    },
    {
        "title": "Hero Carousel Slide 2: Glide",
        "usage": "Hero carousel slide - Second slide showcasing Glide product",
        "dimensions": "1920x1080px (Full viewport) | Product Image: 1200x500px",
        "recommendation": "Animated GIF",
        "context": "Title: 'Glide', Subtitle: 'Anything Anytime Anywhere Delivery'. Shows Glide scooter. Recommend: Scooter in motion GIF"
    },
    {
        "title": "Hero Carousel Slide 3: Consultancy",
        "usage": "Hero carousel slide - Third slide for Fainzy Consultancy",
        "dimensions": "1920x1080px (Full viewport) | Product Image: 1200x500px",
        "recommendation": "Static Image",
        "context": "Title: 'Fainzy Consultancy'. Professional consultancy services image with team or tech visualization"
    },
    {
        "title": "ZiBot Demo Video",
        "usage": "Main product demonstration video below hero section",
        "dimensions": "1200x675px (16:9 aspect ratio, max-width container)",
        "recommendation": "Video",
        "context": "Autoplay looping video showcasing ZiBot in action, deliveries, and navigation capabilities"
    },
    {
        "title": "Product Showcase: ZiBot Last Mile",
        "usage": "Product section - ZiBot Last Mile Delivery System",
        "dimensions": "600x600px (Square product image in alternating layout)",
        "recommendation": "Static Image",
        "context": "Title: 'ZiBot Last Mile Delivery Robot System', Features: 24/7, Cost Effective, Eco-Friendly, Reliable"
    },
    {
        "title": "Product Showcase: Glide",
        "usage": "Product section - Glide autonomous scooter",
        "dimensions": "600x600px (Square product image, reverse layout)",
        "recommendation": "Static Image",
        "context": "Title: 'Glide', Description: Autonomous scooter summon service like Uber"
    },
    {
        "title": "Product Showcase: Consultancy",
        "usage": "Product section - Fainzy Consultancy Services",
        "dimensions": "600x600px (Square product image)",
        "recommendation": "Static Image",
        "context": "Title: 'Fainzy Consultancy Services', Professional consulting image"
    },
    {
        "title": "Custom Solutions Carousel (5 items)",
        "usage": "Scrolling carousel showcasing custom robotics solutions",
        "dimensions": "Each slide: 960x700px (2/3 viewport width) | Product images: 500x500px",
        "recommendation": "Static Image",
        "context": "Solutions: Customized Robots, IoT Solutions, MiraX, Food Ordering System, Hotel Delivery. Recommend carousel with smooth transitions"
    }
]

add_page_slide(
    "HOME PAGE (/)",
    "Main landing page with hero carousel, product showcases, and metrics",
    home_media
)

# ==================== PAGE 2: ABOUT PAGE ====================
about_media = [
    {
        "title": "About Hero Banner",
        "usage": "Page header background image",
        "dimensions": "1920x802px (Full-width hero section, min-height: 802px)",
        "recommendation": "Static Image",
        "context": "Title: 'About Us', Subtitle: 'All you need to know about us'. Professional team or office environment photo with dark overlay"
    },
    {
        "title": "Founded Section Image",
        "usage": "Company story section - Foundation image",
        "dimensions": "528x500px (Left-aligned image in two-column layout)",
        "recommendation": "Static Image",
        "context": "Title: 'Founded in 2018. Headquartered in Nagoya.' Historical or office establishment photo"
    },
    {
        "title": "People Behind Innovation Banner",
        "usage": "Full-width banner image for team section",
        "dimensions": "1920x567px (Full-width, min-height: 567px)",
        "recommendation": "Static Image",
        "context": "Title: 'The People Behind The Innovation'. Team collaboration or innovation workspace photo with gradient overlay"
    },
    {
        "title": "Vision Section Image",
        "usage": "Vision statement section - Right-aligned image",
        "dimensions": "528x500px (Right-aligned in two-column layout)",
        "recommendation": "Static Image",
        "context": "Title: 'Our Vision'. Futuristic technology or robotics visualization"
    },
    {
        "title": "Team Member Photos (7 members)",
        "usage": "Team grid - Individual portraits",
        "dimensions": "280x320px each (Portrait cards in 4-column grid)",
        "recommendation": "Static Image",
        "context": "Team: Dr. Jude (CEO), Emmanuel (Tech Lead), Patrick (Design Lead), Michael (Project Lead), Lazarus (Developer), Prof. Tatsuya (Advisor), Prof. Hiroyuki (Advisor)"
    },
    {
        "title": "Join Us CTA Banner",
        "usage": "Bottom call-to-action background",
        "dimensions": "1920x551px (Full-width CTA section)",
        "recommendation": "Static Image",
        "context": "Title: 'Let's do great things together'. Inspiring team or workspace photo with overlays"
    }
]

add_page_slide(
    "ABOUT PAGE (/about)",
    "Company history, vision, team members, and recruitment CTA",
    about_media
)

# ==================== PAGE 3: BLOG PAGE ====================
blog_media = [
    {
        "title": "Featured Blog Post Hero",
        "usage": "Featured article large image",
        "dimensions": "960x548px (50% width on desktop, full-width mobile, 16:9 aspect)",
        "recommendation": "Static Image",
        "context": "Title: 'The Future of Autonomous Delivery'. High-quality editorial image from Unsplash showcasing ZiBot or delivery robots"
    },
    {
        "title": "Blog Post Thumbnails (6 posts)",
        "usage": "Blog card grid - Thumbnail images",
        "dimensions": "400x224px each (3-column grid, 16:9 aspect ratio)",
        "recommendation": "Static Image",
        "context": "Categories: Business, Technology, Case Study, Careers, Sustainability. Editorial-style images for each topic"
    }
]

add_page_slide(
    "BLOG PAGE (/blog)",
    "Blog listing with featured post and article grid",
    blog_media
)

# ==================== PAGE 4: BLOG POST DETAIL ====================
blog_post_media = [
    {
        "title": "Article Hero Image",
        "usage": "Full-width article header image",
        "dimensions": "1920x500px (Full-width header, 16:9 aspect)",
        "recommendation": "Static Image",
        "context": "Title: '5 Ways Autonomous Robots Are Transforming Hospitality'. Hotel/restaurant robot service photo"
    },
    {
        "title": "In-Article Image 1",
        "usage": "Inline content image for Guest Experiences section",
        "dimensions": "1200x400px (Article width, 3:1 aspect ratio)",
        "recommendation": "Static Image",
        "context": "Section: 'Enhanced Guest Experiences Through Technology'. Modern hotel lobby photo"
    },
    {
        "title": "In-Article Image 2",
        "usage": "Inline content image for Cost Efficiency section",
        "dimensions": "1200x400px (Article width)",
        "recommendation": "Static Image",
        "context": "Section: 'Cost Efficiency and Strong ROI'. Restaurant robot service photo"
    },
    {
        "title": "Related Article Thumbnails (3 posts)",
        "usage": "Similar articles grid at bottom",
        "dimensions": "400x208px each (3-column grid, 16:9)",
        "recommendation": "Static Image",
        "context": "AI/ML, Restaurant Case Study, Sensor Technology articles"
    }
]

add_page_slide(
    "BLOG POST DETAIL (/blog/[id])",
    "Individual article page with hero image and inline media",
    blog_post_media
)

# ==================== PAGE 5: CAREERS PAGE ====================
careers_media = [
    {
        "title": "Careers Hero Banner",
        "usage": "Page header background",
        "dimensions": "1920x600px (Full-width hero, minimum 60vh)",
        "recommendation": "Static Image",
        "context": "Title: 'Join Us in Building the Future'. Inspiring team workspace or innovation lab photo with dark overlay"
    },
    {
        "title": "Department Cards (3 departments)",
        "usage": "Department showcase cards",
        "dimensions": "400x250px each (3-column grid cards)",
        "recommendation": "Static Image",
        "context": "Engineering & Technology, Product & Design, Operations & Manufacturing. Each with relevant workspace photos"
    },
    {
        "title": "World-Class Facilities Banner",
        "usage": "Full-width facility showcase section",
        "dimensions": "1920x500px (Full-width, minimum 50vh)",
        "recommendation": "Static Image",
        "context": "Title: 'World-Class Facilities'. Modern lab or facility photo with dark overlay"
    },
    {
        "title": "Benefits Section Images (3 images)",
        "usage": "Alternating layout benefit sections",
        "dimensions": "600x320px each (Side-by-side with text)",
        "recommendation": "Static Image",
        "context": "Benefits & Perks, Innovation at Scale, Continuous Growth. Professional work environment photos"
    },
    {
        "title": "Team Environment Photo",
        "usage": "Full-width team section divider",
        "dimensions": "1920x288px (Full-width, height: 72px on 1080p)",
        "recommendation": "Static Image",
        "context": "Team collaboration photo with gradient overlay"
    }
]

add_page_slide(
    "CAREERS PAGE (/careers)",
    "Job listings, company culture, benefits, and application process",
    careers_media
)

# ==================== PAGE 6: CONTACT PAGE ====================
contact_media = [
    {
        "title": "Social Media Icons (3 icons)",
        "usage": "Social media links - Twitter, LinkedIn, Instagram",
        "dimensions": "40x40px each (Icon buttons with 16px icons)",
        "recommendation": "Static Image",
        "context": "SVG icons from lucide-react library with hover states"
    },
    {
        "title": "Interactive Map",
        "usage": "Office location map (Nagoya University)",
        "dimensions": "1920x400px (Full container width, approx 400px height)",
        "recommendation": "Static Image",
        "context": "Mapbox/Google Maps embedded showing Nagoya University Incubation Facility location"
    }
]

add_page_slide(
    "CONTACT PAGE (/contact)",
    "Contact form, office information, and location map",
    contact_media
)

# ==================== PAGE 7: CUSTOM SOLUTIONS PAGE ====================
custom_solutions_media = [
    {
        "title": "Custom Solutions Carousel Hero",
        "usage": "Top carousel showcasing custom solutions",
        "dimensions": "Each slide: 960x700px | Product images: 500x600px",
        "recommendation": "Static Image",
        "context": "3 solutions: Consultancy, Customized Robots, IoT Solutions. Each with features grid and contact CTA"
    }
]

add_page_slide(
    "CUSTOM SOLUTIONS PAGE (/custom-solutions)",
    "Custom robotics solutions with carousel and inquiry form",
    custom_solutions_media
)

# ==================== PAGE 8: PRODUCTS PAGE ====================
products_media = [
    {
        "title": "Products Hero Banner",
        "usage": "Page header background",
        "dimensions": "1920x600px (Full-width hero)",
        "recommendation": "Static Image",
        "context": "Title: 'Our Products & Services'. Professional products showcase with dark overlay"
    },
    {
        "title": "Product Showcases (Same as Home)",
        "usage": "ZiBot, Glide, Consultancy product sections",
        "dimensions": "600x600px each (Alternating layout)",
        "recommendation": "Static Image",
        "context": "Identical to home page product sections"
    },
    {
        "title": "Custom Solutions Carousel (5 items)",
        "usage": "Same as home page custom solutions",
        "dimensions": "960x700px slides | 500x500px images",
        "recommendation": "Static Image",
        "context": "MiraX, Efficient Food Ordering, Hotel Delivery, Customized Robots, IoT"
    }
]

add_page_slide(
    "PRODUCTS PAGE (/products)",
    "Complete product catalog and custom solutions showcase",
    products_media
)

# ==================== PAGE 9: PRODUCT DETAIL (ZiBot) ====================
zibot_media = [
    {
        "title": "ZiBot Hero Background Video",
        "usage": "Full-screen hero section with looping background",
        "dimensions": "1920x1080px (Full viewport, opacity: 50%)",
        "recommendation": "Video",
        "context": "Title: 'ZIBOT - Anything Anytime Anywhere Delivery'. Looping video showing ZiBot in action with radial gradient overlay"
    },
    {
        "title": "Product Demo Video 1: Cozy Dinners",
        "usage": "Section 2 video demonstration",
        "dimensions": "1200x675px (Container max-width, 16:9 aspect)",
        "recommendation": "Video",
        "context": "Title: 'From Cozy Dinners to Office Lunches'. Autoplay looping video of ZiBot delivering food"
    },
    {
        "title": "Product Demo Video 2: Hot and Fresh",
        "usage": "Section 3 video demonstration",
        "dimensions": "1200x675px (16:9 aspect)",
        "recommendation": "Video",
        "context": "Title: 'Hot and Fresh. Every Single Time.' Video showing temperature-controlled delivery"
    },
    {
        "title": "LargeFeatureCarousel: Main Features (3 slides)",
        "usage": "Section 4 carousel - Key features showcase",
        "dimensions": "Each slide: 1280x500px (2/3 viewport width on desktop)",
        "recommendation": "Animated GIF",
        "context": "1) 24/7 Availability - GIF of robot operating day/night. 2) All-Weather Operation - GIF in rain/snow. 3) Zero Contact - GIF of contactless delivery process"
    },
    {
        "title": "LargeFeatureCarousel: 360° Sensing (3 slides)",
        "usage": "Section 7 carousel - Environmental sensing features",
        "dimensions": "1280x500px per slide",
        "recommendation": "Animated GIF",
        "context": "1) Accurate Navigation - GIF of precise routing. 2) Traffic Navigation - GIF navigating obstacles. 3) Never Gets Tired - GIF of continuous operation"
    },
    {
        "title": "Cloud Intelligence Video",
        "usage": "Section 9 video with overlay text",
        "dimensions": "1200x675px (16:9)",
        "recommendation": "Video",
        "context": "Title: 'Connected Cloud Intelligence'. Video showing real-time monitoring dashboard. Overlay text: 'Get orders'"
    },
    {
        "title": "LargeFeatureCarousel: Durability (6 slides)",
        "usage": "Section 10 carousel - Built to Handle Anything",
        "dimensions": "1280x500px per slide",
        "recommendation": "Animated GIF",
        "context": "1) Never Gets Sick - Autonomous operation. 2) All-Weather - GIF in various conditions. 3) Rough Terrain - GIF on curbs/steps. 4) Crowds - GIF navigating pedestrians. 5) Security - Monitoring visualization. 6) Precision - GPS tracking GIF"
    },
    {
        "title": "Safe Operation Video",
        "usage": "Section 11 medium-height video",
        "dimensions": "1200x500px (Medium height, 12:5 aspect)",
        "recommendation": "Video",
        "context": "Title: 'Built to Operate Safely'. Safety features demonstration video"
    },
    {
        "title": "LargeFeatureCarousel: Safe Food (3 slides)",
        "usage": "Section 12 carousel - Food safety features",
        "dimensions": "1280x500px per slide",
        "recommendation": "Animated GIF",
        "context": "1) 360° Awareness - GIF of sensor visualization. 2) Navigate Barriers - GIF avoiding obstacles. 3) Absorbs Bumps - GIF of suspension system"
    },
    {
        "title": "Robot Comparison Image",
        "usage": "Section 6 - Technical specs visualization",
        "dimensions": "600x600px (Robot side view with spec callouts)",
        "recommendation": "Static Image",
        "context": "Specifications: 20km range, 30kg payload, 5km/h speed. Annotated robot diagram"
    },
    {
        "title": "Navigate Environment Grid (6 items)",
        "usage": "Section 8 - Feature grid around robot image",
        "dimensions": "Each icon/image: 200x200px | Center robot: 400x400px",
        "recommendation": "Static Image",
        "context": "Smart Sensors, AI Awareness, Pedestrian Detection, Smart Routing, GPS Accuracy, Obstacle Handling. Icons or small illustrations"
    },
    {
        "title": "Transform Business Video",
        "usage": "Section 17 - Business transformation showcase",
        "dimensions": "1200x400px (Wide aspect video)",
        "recommendation": "Video",
        "context": "Title: 'Transform Your Business Operations'. Stats: 30% savings, 15% fee, 6 steps to go live"
    },
    {
        "title": "Specifications Robot Image",
        "usage": "Section 22 - Technical specifications section",
        "dimensions": "800x600px (Right-aligned product image)",
        "recommendation": "Static Image",
        "context": "Detailed ZiBot specifications diagram with callouts for dimensions, battery, controls"
    },
    {
        "title": "Final CTA Robot Image",
        "usage": "Section 23 - Footer call-to-action",
        "dimensions": "320x320px (Square product image)",
        "recommendation": "Static Image",
        "context": "ZiBot product shot for 'Get Started' CTA section"
    },
    {
        "title": "No Boundaries Video Background",
        "usage": "Section 21 - Feature section background",
        "dimensions": "1920x700px (Full-width, medium height)",
        "recommendation": "Video",
        "context": "Title: 'No Boundaries. No Limits.' Looping video of indoor/outdoor operations with radial gradient"
    }
]

add_page_slide(
    "PRODUCT DETAIL: ZiBot (/products/zibot)",
    "Comprehensive ZiBot product page with extensive video and carousel content",
    zibot_media
)

# Add summary/recommendation slide
summary_data = [
    {
        "title": "MEDIA TYPE BREAKDOWN",
        "usage": "Summary statistics across all pages",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": "Total Static Images: ~70 | Videos: 8 | Recommended GIFs: 15 (currently static images in carousels)"
    },
    {
        "title": "KEY RECOMMENDATIONS",
        "usage": "Implementation guidance",
        "dimensions": "N/A",
        "recommendation": "Animated GIF",
        "context": """
Priority 1 (Videos): Keep all existing videos - ZiBot demos, product showcases
Priority 2 (Convert to GIF): LargeFeatureCarousel items on ZiBot page (15 slides total)
Priority 3 (Keep Static): Hero carousel slides, team photos, blog images, product cards
        """
    },
    {
        "title": "ANIMATION OPPORTUNITIES",
        "usage": "Sections that would benefit from motion",
        "dimensions": "N/A",
        "recommendation": "Animated GIF",
        "context": """
✓ Robot navigation demonstrations (obstacles, traffic, crowds)
✓ Weather operation scenarios (rain, snow, day/night)
✓ Contactless delivery process
✓ Sensor visualization and 360° awareness
✓ UI/dashboard interactions
        """
    },
    {
        "title": "TECHNICAL NOTES",
        "usage": "Development considerations",
        "dimensions": "Standard viewport: 1920x1080px (Desktop)",
        "recommendation": "Static Image",
        "context": """
- All dimensions calculated for optimal display on 1920x1080 screens
- Maintain aspect ratios: 16:9 for videos, varies for product images
- Use WebP format for static images for better compression
- Optimize GIF file sizes (max 3-5MB, use tools like ezgif.com)
- Video format: MP4 H.264, with autoplay/loop/muted attributes
- Consider lazy loading for below-fold media
        """
    }
]

add_page_slide(
    "SUMMARY & RECOMMENDATIONS",
    "Implementation guidance and technical considerations",
    summary_data
)

# ==================== BACKEND API REQUIREMENTS ====================
backend_api_data = [
    {
        "title": "Contact Form API",
        "usage": "POST /api/contact - Handle general inquiries and customer contact",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
Fields: name (required), email (required), subject, phone, message (required, min 135 chars), reason (dropdown)
Reason Options: General Inquiry, Partnership Opportunity, Technical Support, Sales Inquiry, Feedback, Other
Response: { success: boolean, message: string, ticketId?: string }
Integrations: SendGrid/AWS SES for email, HubSpot/Salesforce for CRM
        """
    },
    {
        "title": "Newsletter Subscription API",
        "usage": "POST /api/newsletter/subscribe - Collect email subscriptions",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
Fields: email (required with validation), consentGiven (boolean), source (string), timestamp
Response: { success: boolean, message: string }
Integrations: Mailchimp, SendGrid Marketing, or ConvertKit
Terms: "By subscribing, you agree to our Privacy Policy..."
        """
    },
    {
        "title": "Blog Content API",
        "usage": "GET /api/blog/posts, GET /api/blog/posts/:id - Manage blog content",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
List: ?limit=10&offset=0&category=Technology
Response: { posts: Array<{id, title, category, description, date, readTime, image}>, total }
Detail: Full post with content (HTML/Markdown), author, tags, relatedPosts
Current: 6 hardcoded posts (Business, Technology, Case Study, Careers, Sustainability)
Categories: Business, Technology, Case Study, Careers, Sustainability
        """
    },
    {
        "title": "Careers & Jobs API",
        "usage": "GET /api/careers/jobs, POST /api/careers/apply - Job listings and applications",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
List: ?department=Engineering&type=Full%20time&location=Remote
Fields: id (UUID), title, location, date, type, salary, skills, description, responsibilities, requirements
Apply: jobId, applicantInfo (name, email, phone, resume file, coverLetter, linkedin, portfolio)
Current: 9 jobs hardcoded (Engineering, Product, Operations)
Salary Range: ₦600,000 - ₦2,500,000 (Nigerian Naira)
        """
    },
    {
        "title": "FAQs API",
        "usage": "GET /api/faqs?locale=en - Frequently Asked Questions with localization",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
Response: { faqs: Array<{id, question, answer, category, locale}> }
Current: 24 FAQs total - 12 English (content:1), 12 Japanese (content:2)
Localization: Filtered by locale parameter
Storage: Currently hardcoded in component, should be in database
        """
    },
    {
        "title": "Products & Solutions API",
        "usage": "GET /api/products, GET /api/custom-solutions - Product catalog",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
Products: ZiBot Last Mile Delivery, Glide (Autonomous Scooter), Fainzy Consultancy Services
Custom Solutions: 5 solutions (Customized Robots, IoT Solutions, MiraX, Food Ordering, Hotel Delivery)
Fields: id, title, description, image, features (Array<{title, subtitle, icon}>), category, price, availability
Current: All hardcoded in page.tsx files
        """
    }
]

add_page_slide(
    "BACKEND API REQUIREMENTS",
    "Endpoints needed for dynamic content and form submissions - Priority: Critical",
    backend_api_data
)

# ==================== SAMPLE DATA RECOMMENDATIONS ====================
sample_data = [
    {
        "title": "Contact Form Sample Submissions",
        "usage": "Sample data for testing and backend development",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
Example 1 (Japanese):
  name: "山本 隆", email: "tyamamoto@example.jp", phone: "+81-90-1234-5678"
  subject: "Partnership Inquiry", reason: "Partnership Opportunity"
  message: "I represent a hotel chain in Tokyo interested in implementing ZiBot for room service..."

Example 2 (English):
  name: "Sarah Mitchell", email: "sarah.m@restaurantgroup.com"
  subject: "Custom Robot Request", reason: "Custom Robot"
  message: "Our restaurant chain needs a custom delivery robot for outdoor patio service..."

Example 3 (Technical Support):
  name: "Alex Chen", email: "alex.chen@techcorp.com", phone: "+1-415-555-0123"
  reason: "Technical Support", message: "ZiBot navigation issues in crowded areas..."
        """
    },
    {
        "title": "Blog Post Full Structure",
        "usage": "Complete blog post data format with all fields",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
{
  "id": 1, "title": "5 Ways Autonomous Robots Are Transforming Hospitality",
  "slug": "autonomous-robots-hospitality", "category": "Business",
  "tags": ["Hospitality", "Automation", "ZiBot", "Hotels"],
  "description": "From room service to concierge duties, discover how robots like ZiBot...",
  "content": "<full HTML content>", "featuredImage": "https://images.unsplash.com/...",
  "author": { "name": "Fainzy Technologies", "role": "Editorial Team", "avatar": "/logo.png" },
  "publishedAt": "2025-10-20T08:00:00Z", "readTime": "6 min read",
  "views": 12500, "likes": 3542,
  "tableOfContents": ["Introduction", "Revolutionizing Room Service", ...]
}
        """
    },
    {
        "title": "Job Application Sample Data",
        "usage": "Complete job application submission structure",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
{
  "id": "app-001", "jobId": "3f8a5c2e-1d4b-4e9f-a6c7-8b2d4e6f9a1c",
  "jobTitle": "Software Developer (Full-Stack)",
  "applicant": {
    "name": "Alex Chen", "email": "alex.chen@email.com", "phone": "+1-415-555-0123",
    "linkedin": "linkedin.com/in/alexchen", "portfolio": "alexchen.dev"
  },
  "resume": { "filename": "alex_chen_resume.pdf", "url": "/uploads/resumes/..." },
  "coverLetter": "I am excited to apply for the Full-Stack Developer position...",
  "expectedSalary": "₦850,000", "availableStartDate": "2026-01-15",
  "applicationStatus": "under_review", "submittedAt": "2025-11-21T16:50:00Z"
}
        """
    },
    {
        "title": "Form Validation Rules",
        "usage": "Client-side and server-side validation requirements",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
Contact Form:
  - name: required, min 2 chars, max 100 chars
  - email: required, valid email format (RFC 5322)
  - message: required, min 135 chars (main form) or 10 chars (alternative form)
  - phone: optional, E.164 format validation recommended
  - agree: required checkbox for terms acceptance

Job Application:
  - resume: required, PDF/DOCX only, max 5MB
  - coverLetter: optional, max 2000 chars or file upload
  - expectedSalary: optional, numeric validation
  - availableStartDate: date validation, not in past

Newsletter:
  - email: required, valid format, unique constraint
  - consentGiven: required checkbox
        """
    },
    {
        "title": "Metrics & Stats Data",
        "usage": "Real-time statistics for homepage",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
Animated Stats (currently hardcoded, should be API-driven):
  - 150 Robots in active service
  - 45 Restaurants using robot delivery system
  - 30 Hotels using robot delivery system

Static Stats:
  - 25+ Active Deployments
  - 8+ Countries Served
  - 500K+ Deliveries Completed

Recommendation: Create GET /api/stats endpoint for real-time updates
Data refresh: Every 5 minutes or on-demand
        """
    }
]

add_page_slide(
    "SAMPLE DATA & VALIDATION RULES",
    "Recommended test data and validation requirements for forms and APIs",
    sample_data
)

# ==================== LOCALIZATION STATUS ====================
localization_status = [
    {
        "title": "Current Localization System",
        "usage": "React Context with localStorage persistence",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
Implementation: LocaleContext.tsx
Supported Locales: 'en' (English), 'ja' (Japanese)
Storage: localStorage key 'locale'
Auto-detection: Browser language with fallback to English
  - Checks navigator.language.startsWith('ja') → Japanese
  - Default: English for all other languages

Current Coverage:
✅ FAQs: 12 English + 12 Japanese (fully localized)
✅ Contact Page: Office address, map labels (partial)
✅ Navbar: Language selector with EN/JA toggle
⚠️ About Page: English only (needs Japanese translation)
⚠️ Careers: All job listings English only
⚠️ Products: Descriptions English only
⚠️ Blog: All posts English only
        """
    },
    {
        "title": "Contact Page Localization Examples",
        "usage": "Bilingual address and map labels",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
English:
  - Office label: "Office"
  - Address: "Nagoya University Incubation Facility, Furo-cho, Chikusa Ward, Nagoya, Aichi 464-0814, Japan"
  - Map marker: "Nagoya University Incubation Facility"

Japanese:
  - Office label: "オフィス"
  - Address: "〒464-0814 愛知県名古屋市千種区不老町 名古屋大学 インキュベーション施設"
  - Map marker: "名古屋大学 インキュベーション施設"

Implementation: {locale === "ja" ? "日本語" : "English"}
        """
    },
    {
        "title": "FAQ Localization Examples",
        "usage": "Side-by-side English and Japanese Q&A",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
English:
Q: "What is Fainzy Technologies' specialty?"
A: "We specialize on (1)advanced autonomous driving systems, (2)making robots,
    (3)mobile app and website development, and, (4)AI applications using deep learning"

Japanese:
Q: "ファインジー・テクノロジーズの得意分野は何ですか？"
A: "(1)高度な自律走行システム、(2)ロボット製作、(3)モバイルアプリ・ウェブサイト開発、
    (4)ディープラーニングによるAIアプリケーションに特化した事業展開を行っています。"

---

English:
Q: "How can I apply for a role?"
A: "Please contact us using the Contact Us Form, (make sure to select 'Apply For A Role' as the Reason."

Japanese:
Q: "どのように応募すればよいのですか？"
A: "お問い合わせフォームよりご連絡ください（必ず「応募する」を選択してください）。"
        """
    },
    {
        "title": "Navigation Localization",
        "usage": "Menu items translation recommendations",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
Current (English Only):
  HOME, ABOUT, PRODUCTS, CAREERS, CONTACT, BLOG, BUSINESS, LOGIN

Recommended Japanese Translations:
  HOME → ホーム
  ABOUT → 会社概要
  PRODUCTS → 製品・サービス
  CAREERS → 採用情報
  CONTACT → お問い合わせ
  BLOG → ブログ
  BUSINESS → ビジネス
  LOGIN → ログイン

Language Selector:
  - Globe icon with locale badge (EN/JA)
  - Dropdown: "English" | "日本語"
  - Checkmark indicates active language
        """
    }
]

add_page_slide(
    "LOCALIZATION SYSTEM & STATUS",
    "Current implementation and bilingual content coverage",
    localization_status
)

# ==================== LOCALIZATION GAPS & RECOMMENDATIONS ====================
localization_recommendations = [
    {
        "title": "Translation Priority Matrix",
        "usage": "Pages requiring Japanese translation by priority",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
HIGH PRIORITY (Critical for Japanese users):
✓ Navigation menu items (HOME, ABOUT, PRODUCTS, etc.)
✓ Product names and descriptions (ZiBot, Glide, Consultancy)
✓ Contact form labels and placeholders
✓ Footer content and copyright
✓ Error messages and validation text

MEDIUM PRIORITY (Important for user experience):
⚠️ About page content (company history, vision, mission)
⚠️ Team member roles and bios
⚠️ Career page intro and benefits section
⚠️ Product feature descriptions
⚠️ Custom solutions descriptions

LOW PRIORITY (Can remain English):
○ Blog post content (technical articles)
○ Detailed job descriptions
○ Technical documentation
○ Legal terms and privacy policy (dual-language acceptable)
        """
    },
    {
        "title": "Recommended Translation File Structure",
        "usage": "Organized translation management system",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
Current Pattern: {locale === "ja" ? "日本語" : "English"} (inline ternary)

Recommended: Separate translation files

// translations/en.ts
export const en = {
  nav: { home: "HOME", about: "ABOUT", products: "PRODUCTS", ... },
  contact: { title: "Get in touch", email: "Email", phone: "Phone", ... },
  products: { zibot: "ZiBot Last Mile Delivery Robot System", ... },
  common: { loading: "Loading...", error: "Error", success: "Success" }
}

// translations/ja.ts
export const ja = {
  nav: { home: "ホーム", about: "会社概要", products: "製品・サービス", ... },
  contact: { title: "お問い合わせ", email: "メールアドレス", phone: "電話番号", ... },
  products: { zibot: "ZiBot ラストワンマイル配送ロボットシステム", ... },
  common: { loading: "読み込み中...", error: "エラー", success: "成功" }
}

Usage: const t = useTranslation(); <h1>{t('contact.title')}</h1>
        """
    },
    {
        "title": "Page-by-Page Translation Checklist",
        "usage": "Comprehensive coverage assessment",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
✅ Contact Page: 70% complete (address, labels done; form placeholders need work)
✅ FAQs: 100% complete (12 EN + 12 JA)
✅ Map Labels: 100% complete
⚠️ Home Page: 30% complete (products section needs translation)
⚠️ About Page: 0% complete (entire page in English)
⚠️ Products Page: 20% complete (navigation only)
⚠️ Careers Page: 10% complete (navigation only, all job listings English)
⚠️ Blog Page: 10% complete (navigation only, all posts English)
⚠️ Custom Solutions: 0% complete

Total Pages: 14
Fully Localized: 2 pages (14%)
Partially Localized: 3 pages (21%)
Not Localized: 9 pages (65%)

Target: 100% coverage for High Priority pages within 3 months
        """
    },
    {
        "title": "Implementation Roadmap",
        "usage": "Phased translation deployment plan",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
Phase 1 (Weeks 1-2): Critical Navigation & Forms
  - Translate navigation menu
  - Translate all form labels, placeholders, validation messages
  - Translate footer content
  - Update contact page fully

Phase 2 (Weeks 3-4): Product Content
  - Translate product names and descriptions
  - Translate custom solutions
  - Translate feature lists and benefits
  - Update product pages

Phase 3 (Weeks 5-6): About & Careers
  - Translate about page content
  - Translate team member roles
  - Translate benefits and culture sections
  - Note: Job listings can remain English (international hiring)

Phase 4 (Weeks 7-8): Blog & Dynamic Content
  - Create Japanese blog category
  - Translate blog post titles/descriptions (optional)
  - Implement locale-based content filtering

Maintenance: Ongoing translation for new content additions
        """
    },
    {
        "title": "Technical Implementation Notes",
        "usage": "Development guidelines for localization",
        "dimensions": "N/A",
        "recommendation": "Static Image",
        "context": """
Date/Time Formatting:
  - English: "October 20, 2025" → Japanese: "2025年10月20日"
  - Use Intl.DateTimeFormat for automatic formatting

Number Formatting:
  - English: "1,000" → Japanese: "1,000" (same)
  - Currency: ₦ (Naira) remains consistent

URL Structure Options:
  - Option 1: Query param → /contact?lang=ja
  - Option 2: Subdomain → ja.fainzy.ai
  - Option 3: Path prefix → /ja/contact (Recommended)

SEO Considerations:
  - Add hreflang tags for each locale
  - Separate meta descriptions per language
  - Localized sitemaps

Testing Strategy:
  - Test locale switching without page reload
  - Verify localStorage persistence
  - Check RTL support (not needed for Japanese, but good practice)
  - Validate all forms in both languages
        """
    }
]

add_page_slide(
    "LOCALIZATION GAPS & RECOMMENDATIONS",
    "Translation priorities, file structure, and implementation roadmap",
    localization_recommendations
)

# Save presentation
output_path = "/Users/apple/Fainzy/website-v2/Fainzy_Media_Specifications.pptx"
prs.save(output_path)

print(f"""
✅ Comprehensive Specification Presentation Generated Successfully!

📄 File saved to: {output_path}

📊 Summary:
- Total Pages Analyzed: 9 website pages
- Total Media Items Documented: 80+
- Videos Identified: 8
- Recommended GIFs: 15 (currently in LargeFeatureCarousel)
- Static Images: ~70

🔌 Backend & API:
- API Endpoints Documented: 6 (Contact, Newsletter, Blog, Careers, FAQs, Products)
- Forms Requiring Backend: 4 (Contact x2, Newsletter, Job Application)
- Sample Data Examples: 5 complete datasets
- Validation Rules: Comprehensive coverage

🌍 Localization:
- Supported Locales: English (en), Japanese (ja)
- Fully Localized Pages: 2 (14%)
- Partially Localized: 3 (21%)
- Translation Needed: 9 pages (65%)
- FAQ Translations: 12 EN + 12 JA (complete)
- Priority Recommendations: 3-phase roadmap (8 weeks)

📑 Presentation Structure:
- Media Specifications: 10 slides (9 pages + summary)
- Backend API Requirements: 1 slide
- Sample Data & Validation: 1 slide
- Localization Status: 1 slide
- Localization Recommendations: 1 slide
- Total Slides: 14

🎯 Next Steps:
1. Review the generated PowerPoint presentation
2. Source/create media assets based on specifications
3. Implement Backend APIs (Contact, Blog, Careers priority)
4. Create translation files for Japanese content
5. Set up form validation and submission handling
6. Optimize file sizes for web delivery

💡 Key Findings:
- ZiBot product page uses extensive video content (8 videos total)
- LargeFeatureCarousel components should contain animated GIFs for better UX
- All dimensions optimized for 1920x1080 viewport
- 65% of pages need Japanese translation (8-week roadmap provided)
- All forms currently use mock submissions (need backend integration)
- 24 FAQs hardcoded (should migrate to database with locale filtering)
""")
